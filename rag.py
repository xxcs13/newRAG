import os
import re
import json
import logging
import time
import tiktoken
from typing import List, Any, Dict, Tuple, Optional, Union, Literal, Sequence
from dataclasses import dataclass, field
from pathlib import Path
from dotenv import load_dotenv
from concurrent.futures import ThreadPoolExecutor
import inspect

import pandas as pd
from pptx import Presentation
from pypdf import PdfReader  # type: ignore
import pdfplumber

from langchain_community.vectorstores import Chroma
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain.schema.document import Document
from langchain_core.messages import BaseMessage
from langchain.text_splitter import RecursiveCharacterTextSplitter

from langgraph.graph import StateGraph, END
from pydantic import BaseModel, Field

# Environment Configuration
load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
assert OPENAI_API_KEY, "Please set OPENAI_API_KEY in your .env file."

enable_telemetry = os.getenv("ENABLE_TELEMETRY", "false").lower() == "true"
if not enable_telemetry:
    os.environ.setdefault("CHROMA_TELEMETRY", "false")
    os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")
    os.environ.setdefault("ANONYMIZED_TELEMETRY", "false")
    print("Telemetry disabled for privacy and stability. Set ENABLE_TELEMETRY=true in .env to enable.")

_log = logging.getLogger(__name__)

# ============================================================================
# PROMPT SYSTEM
# ============================================================================

def build_system_prompt(instruction: str="", example: str="", pydantic_schema: str="") -> str:
    """Build structured system prompt with instruction, schema and example."""
    delimiter = "\n\n---\n\n"
    schema = f"Your answer should be in JSON and strictly follow this schema:\n```\n{pydantic_schema}\n```"
    if example:
        example = delimiter + example.strip()
    if schema:
        schema = delimiter + schema.strip()
    
    return instruction.strip() + schema + example

class RAGAnswerPrompt:
    """Structured prompt system for RAG answer generation."""
    
    instruction = """
You are an advanced RAG answering system that responds in Traditional Chinese.
Answer questions based on provided documents, synthesizing insights across multiple sources.

FINANCIAL DATA PRIORITY:
For money-related questions (revenue, profit, costs, etc.), prioritize Excel documents first as they contain comprehensive numerical data.

QUESTION TYPE ANALYSIS:
- FACTUAL QUERIES: Provide direct, precise answers with exact data
- ANALYTICAL QUERIES: Provide comprehensive analysis with business insights

CORE PRINCIPLES:
1. Read questions carefully and identify all components
2. Examine content in ALL languages for relevance
3. Cross-reference multiple sources for consistency
4. Distinguish between explicit statements vs inferences
5. Acknowledge when information is incomplete or missing

NUMBER FORMATTING REQUIREMENT:
When presenting numerical values, include both the original number and its Chinese representation in parentheses. Examples:
- 150000 → 150000 (15萬)
- 1500000 → 1500000 (150萬) 
- 25000000 → 25000000 (2500萬)
- 1000000000 → 1000000000 (10億)

PRECISION REQUIREMENT FOR FINANCIAL DATA:
Provide EXACT numbers from source documents, not approximations. Do NOT use terms like "approximately" unless:
1. The source document itself uses approximation language, OR
2. The user specifically asks for approximate figures

Always prioritize precision and accuracy when presenting monetary amounts.
"""

    user_prompt = """
Here is the context:
\"\"\"
{context}
\"\"\"

Here is the question:
"{question}"

REMINDERS:
1. Examine ALL content regardless of language to find relevant information
2. For money-related questions, prioritize Excel-type sources
3. Match response complexity to question type
4. Include information only if actually relevant to the question
"""

    class AnswerSchema(BaseModel):
        step_by_step_analysis: str = Field(description="Analysis following the framework")
        reasoning_summary: str = Field(description="Concise synthesis summary highlighting key evidence")
        relevant_sources: List[str] = Field(description="Source IDs containing information directly used")
        confidence_level: Literal["high", "medium", "low"] = Field(description="Confidence assessment")
        final_answer: str = Field(description="Answer in Traditional Chinese with appropriate depth")

    pydantic_schema = re.sub(r"^ {4}", "", inspect.getsource(AnswerSchema), flags=re.MULTILINE)
    system_prompt = build_system_prompt(instruction)
    system_prompt_with_schema = build_system_prompt(instruction, "", pydantic_schema)

# ============================================================================
# DOCUMENT CHUNKING SYSTEM
# ============================================================================

class CrossPageTextSplitter:
    """Enhanced document chunking with cross-page support."""
    
    def __init__(self, chunk_size: int = 300, chunk_overlap: int = 50):
        """Initialize cross-page text splitter."""
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
            model_name="gpt-4o-mini",
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap
        )
    
    def split_document(self, document_data: Dict) -> List[Document]:
        """Split document into chunks with format-specific handling."""
        doc_type = document_data['metainfo']['document_type']
        
        if doc_type == 'excel':
            return self._split_excel_document(document_data)
        elif doc_type == 'pptx':
            return self._split_pptx_document(document_data)
        else:
            return self._split_cross_page_document(document_data)
    
    def _split_cross_page_document(self, document_data: Dict) -> List[Document]:
        """Split document content across page boundaries for better semantic continuity."""
        pages = document_data['content']['pages']
        if not pages:
            return []
        
        combined_text, page_boundaries = self._combine_pages_with_markers(pages)
        if not combined_text.strip():
            return []
        
        text_chunks = self.text_splitter.split_text(combined_text)
        
        documents = []
        for i, chunk in enumerate(text_chunks):
            if not chunk.strip():
                continue
                
            clean_chunk = self._remove_page_markers(chunk)
            chunk_start = combined_text.find(chunk)
            chunk_end = chunk_start + len(chunk)
            page_range = self._get_page_range(chunk_start, chunk_end, page_boundaries)
            
            metadata = {
                "chunk": i + 1,
                "total_chunks": len(text_chunks),
                "content_type": "cross_page_text"
            }
            
            if len(page_range) == 1:
                metadata["page"] = page_range[0]
            else:
                metadata["page"] = page_range[0]
                metadata["page_range"] = ",".join(map(str, page_range))
                metadata["spans_pages"] = True
            
            documents.append(Document(
                page_content=clean_chunk.strip(),
                metadata=metadata
            ))
        
        return documents
    
    def _combine_pages_with_markers(self, pages: List[Dict]) -> Tuple[str, List[Tuple[int, int, int]]]:
        """Combine pages into continuous text with boundary markers."""
        combined_parts = []
        page_boundaries = []
        current_pos = 0
        
        for page_data in pages:
            page_num = page_data['page']
            page_text = page_data.get('text', '').strip()
            
            if not page_text:
                continue
            
            page_marker = f"\n--- PAGE {page_num} ---\n"
            page_content = page_marker + page_text + "\n"
            
            start_pos = current_pos
            end_pos = current_pos + len(page_content)
            
            page_boundaries.append((page_num, start_pos, end_pos))
            combined_parts.append(page_content)
            current_pos = end_pos
        
        return ''.join(combined_parts), page_boundaries
    
    def _remove_page_markers(self, text: str) -> str:
        """Remove page markers from chunk text."""
        return re.sub(r'\n--- PAGE \d+ ---\n', '\n', text)
    
    def _get_page_range(self, chunk_start: int, chunk_end: int, page_boundaries: List[Tuple[int, int, int]]) -> List[int]:
        """Determine which pages a chunk spans based on position."""
        covered_pages = []
        
        for page_num, start_pos, end_pos in page_boundaries:
            if chunk_start < end_pos and chunk_end > start_pos:
                covered_pages.append(page_num)
        
        return sorted(covered_pages) if covered_pages else [1]
    
    def _split_excel_document(self, document_data: Dict) -> List[Document]:
        """Excel-specific chunking: preserve table structure."""
        chunks = []
        
        for page in document_data['content']['pages']:
            page_text = page.get('text', '')
            if not page_text.strip():
                continue
            
            sheet_chunks = self._split_excel_sheet(page_text, page['page'])
            chunks.extend(sheet_chunks)
        
        return chunks
    
    def _split_excel_sheet(self, sheet_text: str, page_num: int) -> List[Document]:
        """Split Excel sheet while preserving table structure."""
        documents = []
        lines = sheet_text.split('\n')
        
        if not lines:
            return documents
        
        sheet_info = []
        data_rows = []
        headers = ""
        
        for line in lines:
            if line.startswith('Sheet:'):
                sheet_info.append(line)
            elif line.startswith('Headers:'):
                headers = line
                sheet_info.append(line)
            elif line.startswith('Row'):
                data_rows.append(line)
        
        context_header = '\n'.join(sheet_info)
        
        row_groups = []
        current_group = []
        current_size = len(context_header)
        
        for row in data_rows:
            row_size = len(row)
            if current_size + row_size > self.chunk_size and current_group:
                row_groups.append(current_group.copy())
                current_group = [row]
                current_size = len(context_header) + row_size
            else:
                current_group.append(row)
                current_size += row_size
        
        if current_group:
            row_groups.append(current_group)
        
        for i, group in enumerate(row_groups):
            chunk_content = context_header + '\n\n' + '\n'.join(group)
            
            documents.append(Document(
                page_content=chunk_content.strip(),
                metadata={
                    "page": page_num,
                    "chunk": i + 1,
                    "total_chunks": len(row_groups),
                    "content_type": "excel_table",
                    "has_headers": bool(headers)
                }
            ))
        
        return documents
    
    def _split_pptx_document(self, document_data: Dict) -> List[Document]:
        """PPTX-specific chunking: preserve slide structure."""
        chunks = []
        
        for page in document_data['content']['pages']:
            page_text = page.get('text', '')
            if not page_text.strip():
                continue
            
            slide_chunks = self._split_pptx_slide(page_text, page['page'])
            chunks.extend(slide_chunks)
        
        return chunks
    
    def _split_pptx_slide(self, slide_text: str, page_num: int) -> List[Document]:
        """Split PPTX slide while preserving object relationships."""
        documents = []
        content_blocks = self._parse_pptx_content_blocks(slide_text)
        
        if not content_blocks:
            return documents
        
        grouped_blocks = self._group_pptx_content(content_blocks)
        
        for i, group in enumerate(grouped_blocks):
            chunk_content = '\n\n'.join(group['content'])
            
            if chunk_content.strip():
                documents.append(Document(
                    page_content=chunk_content.strip(),
                    metadata={
                        "page": page_num,
                        "chunk": i + 1,
                        "total_chunks": len(grouped_blocks),
                        "content_type": "pptx_slide",
                        "object_types": ",".join(group['types']) if group['types'] else ""
                    }
                ))
        
        return documents
    
    def _parse_pptx_content_blocks(self, slide_text: str) -> List[Dict]:
        """Parse PPTX slide into content blocks by type."""
        blocks = []
        lines = slide_text.split('\n\n')
        
        for line_group in lines:
            if not line_group.strip():
                continue
            
            content_type = 'text'
            if line_group.startswith('Table:'):
                content_type = 'table'
            elif line_group.startswith('Chart:'):
                content_type = 'chart'
            elif line_group.startswith('Image:'):
                content_type = 'image'
            elif line_group.startswith('Text Frame:'):
                content_type = 'text_frame'
            elif line_group.startswith('Group:'):
                content_type = 'group'
            
            blocks.append({
                'content': line_group.strip(),
                'type': content_type,
                'size': len(line_group)
            })
        
        return blocks
    
    def _group_pptx_content(self, content_blocks: List[Dict]) -> List[Dict]:
        """Group PPTX content blocks intelligently."""
        groups = []
        current_group = {'content': [], 'types': set(), 'size': 0}
        
        for block in content_blocks:
            if block['type'] in ['table', 'chart'] and current_group['content']:
                groups.append({
                    'content': current_group['content'].copy(),
                    'types': list(current_group['types'])
                })
                current_group = {'content': [], 'types': set(), 'size': 0}
            
            current_group['content'].append(block['content'])
            current_group['types'].add(block['type'])
            current_group['size'] += block['size']
            
            if block['type'] in ['table', 'chart']:
                groups.append({
                    'content': current_group['content'].copy(),
                    'types': list(current_group['types'])
                })
                current_group = {'content': [], 'types': set(), 'size': 0}
            elif current_group['size'] > self.chunk_size:
                groups.append({
                    'content': current_group['content'].copy(),
                    'types': list(current_group['types'])
                })
                current_group = {'content': [], 'types': set(), 'size': 0}
        
        if current_group['content']:
            groups.append({
                'content': current_group['content'],
                'types': list(current_group['types'])
            })
        
        return groups

# ============================================================================
# PARENT PAGE AGGREGATION
# ============================================================================

class EnhancedParentPageAggregator:
    """Enhanced parent page retrieval with support for cross-page chunks."""
    
    def __init__(self, parsed_reports: List[Dict]):
        """Initialize with parsed document reports."""
        self.parsed_reports = parsed_reports
        self.page_content_map = self._build_page_content_map()
    
    def _build_page_content_map(self) -> Dict[int, str]:
        """Build mapping from page numbers to full page content."""
        page_map = {}
        for report in self.parsed_reports:
            for page_data in report['report']['content']['pages']:
                page_num = page_data['page']
                page_map[page_num] = page_data['text']
        return page_map
    
    def aggregate_to_parent_pages(self, chunk_results: List[Dict]) -> List[Dict]:
        """Extract parent pages from chunks with cross-page support."""
        seen_page_combinations = set()
        parent_results = []
        
        for chunk_result in chunk_results:
            page_coverage = self._get_chunk_page_coverage(chunk_result)
            page_combination_key = tuple(sorted(page_coverage))
            
            if page_combination_key not in seen_page_combinations:
                seen_page_combinations.add(page_combination_key)
                combined_content = self._get_combined_page_content(page_coverage)
                
                parent_result = {
                    'text': combined_content,
                    'page': page_coverage[0],
                    'page_range': ",".join(map(str, page_coverage)) if len(page_coverage) > 1 else None,
                    'spans_pages': len(page_coverage) > 1,
                    'distance': chunk_result['distance'],
                    'source_file': chunk_result['source_file'],
                    'document_type': chunk_result['document_type'],
                    'metadata': chunk_result['metadata']
                }
                parent_results.append(parent_result)
        
        return parent_results
    
    def _get_chunk_page_coverage(self, chunk_result: Dict) -> List[int]:
        """Determine which pages a chunk covers."""
        metadata = chunk_result.get('metadata', {})
        
        if metadata.get('spans_pages', False) and 'page_range' in metadata:
            page_range_str = metadata['page_range']
            return [int(page) for page in page_range_str.split(',')]
        else:
            return [chunk_result['page']]
    
    def _get_combined_page_content(self, page_numbers: List[int]) -> str:
        """Combine content from multiple pages."""
        if len(page_numbers) == 1:
            return self.page_content_map.get(page_numbers[0], '')
        
        combined_parts = []
        for page_num in sorted(page_numbers):
            page_content = self.page_content_map.get(page_num, '')
            if page_content.strip():
                combined_parts.append(f"[Page {page_num}]\n{page_content}")
        
        return '\n\n--- PAGE BREAK ---\n\n'.join(combined_parts)

# Aliases for backward compatibility
TextSplitter = CrossPageTextSplitter
ParentPageAggregator = EnhancedParentPageAggregator

# ============================================================================
# WORKFLOW STATE MANAGEMENT
# ============================================================================

@dataclass
class GraphState:
    """State management for the RAG workflow graph."""
    docs: Sequence[Union[str, Document]] = field(default_factory=list)
    vectorstore: Any = None
    question: str = ""
    retrieved_docs: List[Document] = field(default_factory=list)
    answer: str = ""
    
    # Enhanced state fields
    parsed_reports: List[Dict] = field(default_factory=list)
    vector_results: List[Dict] = field(default_factory=list)
    reranked_results: List[Dict] = field(default_factory=list)
    final_context: str = ""
    structured_answer: Dict = field(default_factory=dict)
    skip_parsing: bool = False
    
    # Performance tracking fields
    start_time: float = 0.0
    end_time: float = 0.0
    retrieval_time: float = 0.0
    generation_time: float = 0.0
    total_time: float = 0.0
    input_tokens: int = 0
    output_tokens: int = 0
    throughput_tokens_per_second: float = 0.0

# ============================================================================
# DOCUMENT INGESTION NODES
# ============================================================================

def ingest_node(state: GraphState) -> GraphState:
    """Parse and ingest documents using unified parsing system."""
    print("Starting document ingestion...")
    
    parsed_reports = []
    parser = UnifiedDocumentParser()
    text_splitter = CrossPageTextSplitter()
    
    if state.docs and isinstance(state.docs[0], str):
        successful_count = 0
        failed_count = 0
        
        for file_path in state.docs:
            if os.path.exists(str(file_path)):
                try:
                    report = parser.parse_document(str(file_path))
                    
                    if report['metainfo']['document_type'] == 'failed':
                        print(f"Skipping failed document: {file_path}")
                        failed_count += 1
                        continue
                    
                    chunks = text_splitter.split_document(report)
                    
                    for chunk in chunks:
                        chunk.metadata.update({
                            "source_file": os.path.basename(str(file_path)),
                            "document_type": report['metainfo'].get('document_type', 'unknown'),
                            "sha1_name": report['metainfo'].get('sha1_name', '')
                        })
                    
                    parsed_reports.append({
                        'file_path': file_path,
                        'report': report,
                        'chunks': chunks
                    })
                    
                    successful_count += 1
                    print(f"Successfully parsed: {file_path} ({len(chunks)} chunks)")
                    
                except Exception as e:
                    print(f"Failed to parse {file_path}: {e}")
                    failed_count += 1
            else:
                print(f"Warning: File not found: {file_path}")
                failed_count += 1
        
        print(f"Parsing summary: {successful_count} successful, {failed_count} failed")
        
        if successful_count == 0:
            print("Error: No documents were successfully parsed. Check file formats and paths.")
            raise RuntimeError("Document parsing failed for all files")
    
    all_chunks = []
    for parsed_report in parsed_reports:
        all_chunks.extend(parsed_report['chunks'])
    
    print(f"Ingested {len(parsed_reports)} documents with {len(all_chunks)} chunks")
    
    return GraphState(
        docs=all_chunks,
        vectorstore=state.vectorstore,
        question=state.question,
        parsed_reports=parsed_reports
    )

def embed_node(state: GraphState) -> GraphState:
    """Create vector embeddings for document chunks."""
    print("Creating vector embeddings...")
    
    if state.vectorstore is None and state.docs:
        vs_manager = VectorStoreManager()
        document_list = [doc for doc in state.docs if isinstance(doc, Document)]
        
        if document_list:
            vectorstore = vs_manager.create_vectorstore(document_list)
            print(f"Created vector store with {len(document_list)} documents")
        else:
            vectorstore = None
    else:
        vectorstore = state.vectorstore
    
    return GraphState(
        docs=state.docs,
        vectorstore=vectorstore,
        question=state.question,
        parsed_reports=state.parsed_reports
    )

# ============================================================================
# PERFORMANCE UTILITIES
# ============================================================================

def count_tokens(text: str, model: str = "gpt-4o-mini") -> int:
    """Count tokens in text using OpenAI's official tiktoken library."""
    try:
        encoding = tiktoken.encoding_for_model(model)
        return len(encoding.encode(text))
    except KeyError:
        try:
            encoding = tiktoken.get_encoding("o200k_base")
            return len(encoding.encode(text))
        except Exception:
            encoding = tiktoken.get_encoding("cl100k_base")
            return len(encoding.encode(text))

def calculate_throughput(tokens: int, time_seconds: float) -> float:
    """Calculate tokens per second throughput."""
    if time_seconds <= 0:
        return 0.0
    return round(tokens / time_seconds, 2)

# ============================================================================
# FILE INPUT SYSTEM
# ============================================================================

def get_user_files() -> List[str]:
    """Get file paths from user input."""
    print("Please enter file paths for processing.")
    print("Supported formats: PDF, PPTX, XLS/XLSX")
    print("Enter file paths one by one. Type 'done' when finished...")
    
    files = []
    while True:
        file_path = input(f"File {len(files) + 1} (or 'done'): ").strip()
        
        if file_path.lower() == 'done':
            break
        
        if not file_path:
            continue
        
        if os.path.exists(file_path):
            file_ext = Path(file_path).suffix.lower()
            if file_ext in ['.pdf', '.pptx', '.ppt', '.xls', '.xlsx']:
                files.append(file_path)
                print(f"Added: {file_path}")
            else:
                print(f"Unsupported file type: {file_ext}. Supported: .pdf, .pptx, .ppt, .xls, .xlsx")
        else:
            print(f"File not found: {file_path}")
    
    return files

# ============================================================================
# VECTOR DATABASE MANAGEMENT
# ============================================================================

class VectorStoreManager:
    """Vector database persistence and loading."""
    
    def __init__(self, persist_directory: str = "chromadb_v1"):
        self.persist_directory = persist_directory
        self.embeddings = OpenAIEmbeddings(model="text-embedding-3-small")
    
    def vectorstore_exists(self) -> bool:
        """Check if vector database exists in persist directory."""
        if not os.path.exists(self.persist_directory):
            return False
        
        required_files = ['chroma.sqlite3']
        for file in required_files:
            if not os.path.exists(os.path.join(self.persist_directory, file)):
                return False
        
        return True
    
    def load_existing_vectorstore(self) -> Optional[Chroma]:
        """Load existing vector database from persist directory."""
        try:
            if not self.vectorstore_exists():
                return None
            
            vectorstore = Chroma(
                persist_directory=self.persist_directory,
                embedding_function=self.embeddings
            )
            
            test_results = vectorstore.similarity_search("test", k=1)
            print(f"Loaded existing vector database with {vectorstore._collection.count()} documents")
            return vectorstore
            
        except Exception as e:
            print(f"Error loading existing vector database: {e}")
            return None
    
    def create_vectorstore(self, documents: List[Document]) -> Chroma:
        """Create new vector database from documents."""
        vectorstore = Chroma.from_documents(
            documents,
            embedding=self.embeddings,
            persist_directory=self.persist_directory
        )
        print(f"Created new vector database with {len(documents)} documents")
        return vectorstore
    
    def get_vectorstore_stats(self, vectorstore: Chroma) -> Dict[str, Any]:
        """Get statistics about the vector database."""
        try:
            count = vectorstore._collection.count()
            return {
                "document_count": count,
                "persist_directory": self.persist_directory,
                "embedding_model": "text-embedding-3-small"
            }
        except Exception as e:
            return {"error": str(e)}

# ============================================================================
# DOCUMENT PARSING SYSTEM
# ============================================================================

class PDFParser:
    """PDF text extraction with number formatting correction."""
    
    def parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """Extract text and tables from PDF file."""
        try:
            filename = Path(file_path).name
            pages = []
            
            try:
                pages = self._extract_with_pdfplumber(file_path)
                print(f"Successfully parsed PDF with pdfplumber: {filename}")
            except Exception as e:
                print(f"pdfplumber failed, trying pypdf: {e}")
                try:
                    pages = self._extract_with_pypdf(file_path)
                    print(f"Successfully parsed PDF with pypdf: {filename}")
                except Exception as e2:
                    print(f"Both PDF extraction methods failed: {e2}")
                    return self._create_fallback_report(file_path)
            
            processed_pages = []
            for page_data in pages:
                processed_text = self._post_process_text(page_data['text'])
                if processed_text.strip():
                    processed_pages.append({
                        'page': page_data['page'],
                        'text': processed_text.strip()
                    })
            
            report = {
                'metainfo': {
                    'sha1_name': filename.rsplit('.', 1)[0],
                    'filename': filename,
                    'pages_amount': len(processed_pages),
                    'text_blocks_amount': len(processed_pages),
                    'tables_amount': 0,
                    'pictures_amount': 0,
                    'document_type': 'pdf'
                },
                'content': {'pages': processed_pages},
                'tables': [],
                'pictures': []
            }
            
            print(f"Successfully parsed PDF: {filename} ({len(processed_pages)} pages)")
            return report
            
        except Exception as e:
            print(f"Error parsing PDF file {file_path}: {e}")
            return self._create_fallback_report(file_path)
    
    def _extract_with_pdfplumber(self, file_path: str) -> List[Dict]:
        """Extract text and tables using pdfplumber."""
        pages = []
        
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text_parts = []
                
                standard_text = page.extract_text()
                if standard_text:
                    text_parts.append(standard_text)
                
                try:
                    layout_text = page.extract_text(layout=True, x_tolerance=1, y_tolerance=1)
                    if layout_text and layout_text != standard_text:
                        text_parts.append("=== Layout Preserved ===")
                        text_parts.append(layout_text)
                except:
                    pass
                
                try:
                    tables = page.extract_tables()
                    if tables:
                        text_parts.append("=== Tables ===")
                        for i, table in enumerate(tables):
                            if table:
                                table_text = self._format_table_text(table)
                                text_parts.append(f"Table {i+1}:\n{table_text}")
                except:
                    pass
                
                combined_text = '\n\n'.join(text_parts)
                if combined_text.strip():
                    pages.append({
                        'page': page_num,
                        'text': combined_text.strip()
                    })
        
        return pages
    
    def _extract_with_pypdf(self, file_path: str) -> List[Dict]:
        """Extract text using pypdf as fallback."""
        pages = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    pages.append({
                        'page': page_num,
                        'text': text.strip()
                    })
        
        return pages
    
    def _format_table_text(self, table: List[List]) -> str:
        """Format extracted table data into readable text."""
        if not table:
            return ""
        
        formatted_rows = []
        for row in table:
            if row:
                clean_cells = []
                for cell in row:
                    if cell is not None:
                        cell_text = str(cell).strip()
                        cell_text = self._fix_number_formatting(cell_text)
                        clean_cells.append(cell_text)
                    else:
                        clean_cells.append("")
                formatted_rows.append(" | ".join(clean_cells))
        
        return "\n".join(formatted_rows)
    
    def _post_process_text(self, text: str) -> str:
        """Fix number formatting and spacing issues in extracted text."""
        if not text:
            return text
        
        processed_text = self._fix_number_formatting(text)
        processed_text = self._fix_spacing_issues(processed_text)
        processed_text = self._fix_symbols(processed_text)
        
        return processed_text
    
    def _fix_number_formatting(self, text: str) -> str:
        """Correct decimal points, thousand separators, and currency symbols."""
        text = re.sub(r'(\d)\s*[,，]\s*(\d{3})', r'\1,\2', text)
        text = re.sub(r'(\d)\s*[.．]\s*(\d)', r'\1.\2', text)
        text = re.sub(r'(\d)\s*[oO]\s*(\d)', r'\1.0\2', text)
        text = re.sub(r'(\d)\s+(\d)', r'\1\2', text)
        text = re.sub(r'(\d)\s*[%％]', r'\1%', text)
        text = re.sub(r'[$＄]\s*(\d)', r'$\1', text)
        text = re.sub(r'([NT$]+)\s*(\d)', r'\1\2', text)
        text = re.sub(r'[-－—]\s*(\d)', r'-\1', text)
        
        return text
    
    def _fix_spacing_issues(self, text: str) -> str:
        """Fix spacing and line break issues."""
        text = re.sub(r'(\w+)-\s*\n\s*(\w+)', r'\1\2', text)
        text = re.sub(r' {2,}', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        return text
    
    def _fix_symbols(self, text: str) -> str:
        """Fix currency and other symbols."""
        text = re.sub(r'[＄$]', '$', text)
        text = re.sub(r'[％%]', '%', text)
        text = re.sub(r'[（(]', '(', text)
        text = re.sub(r'[）)]', ')', text)
        
        return text
    
    def _create_fallback_report(self, file_path: str) -> Dict[str, Any]:
        """Create minimal report when PDF parsing fails."""
        filename = Path(file_path).name
        return {
            'metainfo': {
                'sha1_name': filename.rsplit('.', 1)[0],
                'filename': filename,
                'pages_amount': 0,
                'text_blocks_amount': 0,
                'tables_amount': 0,
                'pictures_amount': 0,
                'document_type': 'failed'
            },
            'content': {'pages': []},
            'tables': [],
            'pictures': []
        }

class PPTXParser:
    """PPTX content extraction including tables, charts, and images."""
    
    def parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """Parse PPTX file for all types of content."""
        try:
            filename = Path(file_path).name
            prs = Presentation(file_path)
            
            pages = []
            tables_found = 0
            charts_found = 0
            images_found = 0
            other_objects_found = 0
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = self._extract_slide_content(slide)
                
                tables_found += slide_content['stats']['tables']
                charts_found += slide_content['stats']['charts'] 
                images_found += slide_content['stats']['images']
                other_objects_found += slide_content['stats']['other_objects']
                
                if slide_content['combined_text'].strip():
                    pages.append({
                        'page': slide_num,
                        'text': slide_content['combined_text'].strip()
                    })
            
            report = {
                'metainfo': {
                    'sha1_name': filename.rsplit('.', 1)[0],
                    'filename': filename,
                    'pages_amount': len(pages),
                    'text_blocks_amount': len(pages),
                    'tables_amount': tables_found,
                    'pictures_amount': images_found,
                    'document_type': 'pptx'
                },
                'content': {'pages': pages},
                'tables': [],
                'pictures': []
            }
            
            print(f"Successfully parsed PPTX: {filename}")
            print(f"  - {len(pages)} slides with content")
            print(f"  - {tables_found} tables, {charts_found} charts")
            print(f"  - {images_found} images, {other_objects_found} other objects")
            return report
            
        except Exception as e:
            print(f"Error parsing PPTX file {file_path}: {e}")
            return self._create_fallback_report(file_path)
    
    def _extract_slide_content(self, slide) -> Dict[str, Any]:
        """Extract all types of content from a single slide."""
        content_parts = []
        stats = {'tables': 0, 'charts': 0, 'images': 0, 'other_objects': 0}
        
        if hasattr(slide, 'shapes'):
            for shape in slide.shapes:
                try:
                    shape_content = self._process_shape(shape, stats)
                    if shape_content:
                        content_parts.append(shape_content)
                except Exception as e:
                    print(f"Warning: Error processing shape: {e}")
                    stats['other_objects'] += 1
                    continue
        
        return {
            'combined_text': '\n\n'.join(content_parts),
            'stats': stats
        }
    
    def _process_shape(self, shape, stats: Dict[str, int]) -> str:
        """Process individual shape and extract relevant content."""
        content_parts = []
        
        try:
            if self._has_table(shape):
                try:
                    table_text = self._extract_table_text(shape.table)
                    if table_text:
                        content_parts.append(f"Table:\n{table_text}")
                        stats['tables'] += 1
                except Exception:
                    pass
            
            elif self._has_chart(shape):
                try:
                    chart_text = self._extract_chart_text(shape.chart)
                    if chart_text:
                        content_parts.append(f"Chart:\n{chart_text}")
                        stats['charts'] += 1
                except Exception:
                    pass
            
            elif self._is_image_shape(shape):
                image_info = self._extract_image_info(shape)
                if image_info:
                    content_parts.append(f"Image: {image_info}")
                    stats['images'] += 1
            
            elif self._is_group_shape(shape):
                group_content = self._extract_group_content(shape, stats)
                if group_content:
                    content_parts.append(f"Group: {group_content}")
            
            elif hasattr(shape, 'text_frame') and shape.text_frame:
                text_frame_content = self._extract_text_frame(shape.text_frame)
                if text_frame_content:
                    content_parts.append(f"Text Frame: {text_frame_content}")
            
            elif self._is_special_shape(shape):
                special_content = self._extract_special_shape_content(shape)
                if special_content:
                    content_parts.append(f"Special Object: {special_content}")
                    stats['other_objects'] += 1
            
            elif hasattr(shape, 'text') and shape.text and shape.text.strip():
                content_parts.append(f"Text: {shape.text.strip()}")
            
            else:
                other_content = self._extract_other_shape_content(shape)
                if other_content:
                    content_parts.append(f"Shape: {other_content}")
                    stats['other_objects'] += 1
        
        except Exception as e:
            print(f"Warning: Error in shape processing: {e}")
            stats['other_objects'] += 1
        
        return '\n'.join(content_parts) if content_parts else ""
    
    def _is_image_shape(self, shape) -> bool:
        """Check if shape is an image/picture."""
        try:
            if hasattr(shape, 'image'):
                return True
            
            if hasattr(shape, 'shape_type'):
                PICTURE_TYPE = 13
                return shape.shape_type == PICTURE_TYPE
            
            return False
        except Exception:
            return False
    
    def _is_group_shape(self, shape) -> bool:
        """Check if shape is a group."""
        try:
            if hasattr(shape, 'shape_type'):
                GROUP_TYPE = 6
                return shape.shape_type == GROUP_TYPE
            return False
        except Exception:
            return False
    
    def _is_special_shape(self, shape) -> bool:
        """Check if shape is SmartArt or other special object."""
        try:
            if hasattr(shape, 'shape_type'):
                SMART_ART_TYPE = 15
                MEDIA_TYPE = 16
                OLE_OBJECT_TYPE = 7
                
                return shape.shape_type in [SMART_ART_TYPE, MEDIA_TYPE, OLE_OBJECT_TYPE]
            return False
        except Exception:
            return False
    
    def _has_table(self, shape) -> bool:
        """Safely check if shape contains a table."""
        try:
            return hasattr(shape, 'table') and shape.table is not None
        except Exception:
            return False
    
    def _has_chart(self, shape) -> bool:
        """Safely check if shape contains a chart."""
        try:
            return hasattr(shape, 'chart') and shape.chart is not None
        except Exception:
            return False
    
    def _extract_special_shape_content(self, shape) -> str:
        """Extract content from special shapes like SmartArt."""
        try:
            content_parts = []
            
            if hasattr(shape, 'text') and shape.text:
                content_parts.append(f"Text: {shape.text.strip()}")
            
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text_frame_content = self._extract_text_frame(shape.text_frame)
                if text_frame_content:
                    content_parts.append(f"Content: {text_frame_content}")
            
            if hasattr(shape, 'shape_type'):
                content_parts.append(f"Type: {shape.shape_type}")
            
            return ', '.join(content_parts) if content_parts else "Special object"
        except Exception as e:
            print(f"Warning: Error extracting special shape: {e}")
            return "Special object"
    
    def _extract_text_frame(self, text_frame) -> str:
        """Extract text from text frame."""
        try:
            if hasattr(text_frame, 'text') and text_frame.text:
                return text_frame.text.strip()
            elif hasattr(text_frame, 'paragraphs'):
                paragraphs = []
                for paragraph in text_frame.paragraphs:
                    if hasattr(paragraph, 'text') and paragraph.text:
                        paragraphs.append(paragraph.text.strip())
                return '\n'.join(paragraphs) if paragraphs else ""
        except Exception as e:
            print(f"Warning: Error extracting text frame: {e}")
        return ""
    
    def _extract_image_info(self, shape) -> str:
        """Extract basic information about images."""
        try:
            info_parts = []
            if hasattr(shape, 'name') and shape.name:
                info_parts.append(f"Name: {shape.name}")
            
            if hasattr(shape, 'element') and hasattr(shape.element, 'get'):
                alt_text = shape.element.get('alt', '')
                if alt_text:
                    info_parts.append(f"Alt text: {alt_text}")
            
            return ', '.join(info_parts) if info_parts else "Image object"
        except Exception as e:
            print(f"Warning: Error extracting image info: {e}")
            return "Image object"
    
    def _extract_group_content(self, group_shape, stats: Dict[str, int]) -> str:
        """Extract content from grouped shapes."""
        try:
            group_parts = []
            if hasattr(group_shape, 'shapes'):
                for shape in group_shape.shapes:
                    shape_content = self._process_shape(shape, stats)
                    if shape_content:
                        group_parts.append(shape_content)
            return '\n'.join(group_parts) if group_parts else ""
        except Exception as e:
            print(f"Warning: Error extracting group content: {e}")
            return "Group object"
    
    def _extract_other_shape_content(self, shape) -> str:
        """Extract content from other shape types."""
        try:
            content_parts = []
            
            if hasattr(shape, 'name') and shape.name:
                content_parts.append(f"Shape name: {shape.name}")
            
            if hasattr(shape, 'text') and shape.text and shape.text.strip():
                content_parts.append(f"Text: {shape.text.strip()}")
            
            if hasattr(shape, 'shape_type'):
                content_parts.append(f"Type: {shape.shape_type}")
            
            return ', '.join(content_parts) if content_parts else ""
        except Exception as e:
            print(f"Warning: Error extracting other shape content: {e}")
            return ""
    
    def _extract_table_text(self, table) -> str:
        """Extract text content from PPTX table."""
        try:
            table_rows = []
            for row in table.rows:
                row_cells = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    row_cells.append(cell_text)
                table_rows.append(" | ".join(row_cells))
            return "\n".join(table_rows)
        except Exception as e:
            print(f"Error extracting table text: {e}")
            return ""
    
    def _extract_chart_text(self, chart) -> str:
        """Extract text content from PPTX chart."""
        try:
            chart_parts = []
            
            if hasattr(chart, 'chart_title') and chart.chart_title and hasattr(chart.chart_title, 'text_frame'):
                try:
                    chart_parts.append(f"Title: {chart.chart_title.text_frame.text}")
                except:
                    pass
            
            if hasattr(chart, 'plots'):
                for plot in chart.plots:
                    if hasattr(plot, 'series'):
                        for series in plot.series:
                            try:
                                if hasattr(series, 'name') and series.name:
                                    chart_parts.append(f"Series: {series.name}")
                            except:
                                pass
            
            return "\n".join(chart_parts) if chart_parts else ""
        except Exception as e:
            print(f"Error extracting chart text: {e}")
            return ""
    
    def _create_fallback_report(self, file_path: str) -> Dict[str, Any]:
        """Create minimal report when PPTX parsing fails."""
        filename = Path(file_path).name
        return {
            'metainfo': {
                'sha1_name': filename.rsplit('.', 1)[0],
                'filename': filename,
                'pages_amount': 0,
                'text_blocks_amount': 0,
                'tables_amount': 0,
                'pictures_amount': 0,
                'document_type': 'failed'
            },
            'content': {'pages': []},
            'tables': [],
            'pictures': []
        }

class ExcelParser:
    """Excel file parsing for all sheets and data."""
    
    def parse_excel(self, file_path: str) -> Dict[str, Any]:
        """Parse Excel file and return structured report."""
        try:
            file_path_obj = Path(file_path)
            filename = file_path_obj.name
            
            if file_path_obj.suffix.lower() == '.xls':
                excel_data = pd.read_excel(file_path, sheet_name=None, engine='xlrd')
            else:
                excel_data = pd.read_excel(file_path, sheet_name=None, engine='openpyxl')
            
            pages = []
            page_num = 1
            total_text_blocks = 0
            
            for sheet_name, df in excel_data.items():
                if df.empty:
                    continue
                
                sheet_text = f"Sheet: {sheet_name}\n\n"
                
                if not df.columns.empty:
                    headers = " | ".join([str(col) for col in df.columns])
                    sheet_text += f"Headers: {headers}\n\n"
                
                for row_idx, (idx, row) in enumerate(df.iterrows(), start=1):
                    row_text = " | ".join([str(val) if pd.notna(val) else "" for val in row.values])
                    sheet_text += f"Row {row_idx}: {row_text}\n"
                
                pages.append({
                    'page': page_num,
                    'text': sheet_text.strip()
                })
                
                page_num += 1
                total_text_blocks += len(df)
            
            report = {
                'metainfo': {
                    'sha1_name': filename.rsplit('.', 1)[0],
                    'filename': filename,
                    'pages_amount': len(pages),
                    'text_blocks_amount': total_text_blocks,
                    'tables_amount': len(excel_data),
                    'pictures_amount': 0,
                    'document_type': 'excel'
                },
                'content': {'pages': pages},
                'tables': [],
                'pictures': []
            }
            
            print(f"Successfully parsed Excel file: {filename} ({len(pages)} sheets)")
            return report
            
        except Exception as e:
            print(f"Error parsing Excel file {file_path}: {e}")
            return self._create_fallback_report(file_path)
    
    def _create_fallback_report(self, file_path: str) -> Dict[str, Any]:
        """Create minimal report when Excel parsing fails."""
        filename = Path(file_path).name
        return {
            'metainfo': {
                'sha1_name': filename.rsplit('.', 1)[0],
                'filename': filename,
                'pages_amount': 0,
                'text_blocks_amount': 0,
                'tables_amount': 0,
                'pictures_amount': 0,
                'document_type': 'failed'
            },
            'content': {'pages': []},
            'tables': [],
            'pictures': []
        }

class UnifiedDocumentParser:
    """Route documents to appropriate parser by file extension."""
    
    def __init__(self):
        self.pdf_parser = PDFParser()
        self.pptx_parser = PPTXParser()
        self.excel_parser = ExcelParser()
    
    def parse_document(self, file_path: str) -> Dict[str, Any]:
        """Parse document using appropriate parser based on file extension."""
        try:
            file_path_obj = Path(file_path)
            file_extension = file_path_obj.suffix.lower()
            
            if file_extension == '.pdf':
                print(f"Parsing PDF: {file_path}")
                return self.pdf_parser.parse_pdf(file_path)
            elif file_extension in ['.pptx', '.ppt']:
                print(f"Parsing PPTX: {file_path}")
                return self.pptx_parser.parse_pptx(file_path)
            elif file_extension in ['.xls', '.xlsx']:
                print(f"Parsing Excel: {file_path}")
                return self.excel_parser.parse_excel(file_path)
            else:
                print(f"Unsupported file type: {file_extension}")
                return self._create_fallback_report(file_path)
            
        except Exception as e:
            print(f"Error parsing document {file_path}: {e}")
            return self._create_fallback_report(file_path)
    
    def _create_fallback_report(self, file_path: str) -> Dict[str, Any]:
        """Create minimal report when document parsing fails."""
        filename = Path(file_path).name
        return {
            'metainfo': {
                'sha1_name': filename.rsplit('.', 1)[0],
                'filename': filename,
                'pages_amount': 0,
                'text_blocks_amount': 0,
                'tables_amount': 0,
                'pictures_amount': 0,
                'document_type': 'failed'
            },
            'content': {'pages': []},
            'tables': [],
            'pictures': []
        }

# ============================================================================
# RETRIEVAL SYSTEM
# ============================================================================

class RetrievalRankingSingleBlock(BaseModel):
    """Rank retrieved text block relevance to a query."""
    reasoning: str = Field(description="Analysis of the block and how it relates to the query")
    relevance_score: float = Field(description="Relevance score from 0 to 1")

class RetrievalRankingMultipleBlocks(BaseModel):
    """Rank retrieved multiple text blocks relevance to a query."""
    block_rankings: List[RetrievalRankingSingleBlock] = Field(
        description="A list of text blocks and their associated relevance scores."
    )

class LLMReranker:
    """Enhanced LLM-based reranking with detailed reasoning for page relevance scoring."""
    
    def __init__(self):
        self.llm = ChatOpenAI(model="gpt-4o-mini", temperature=0)
        
        self.system_prompt_multiple = """
You are a RAG retrieval ranker.

You will receive a query and several retrieved text blocks. Your task is to evaluate and score each block based on its relevance to the query.

Instructions:

1. Reasoning: 
   Analyze each block by identifying key information and how it relates to the query. Consider whether the block provides direct answers, partial insights, or background context. Explain your reasoning in a few sentences, referencing specific elements to justify your evaluation.

2. Relevance Score (0 to 1, in increments of 0.1):
   0 = Completely Irrelevant: No connection to the query
   0.1 = Virtually Irrelevant: Very slight connection
   0.2 = Very Slightly Relevant: Minimal connection
   0.3 = Slightly Relevant: Small aspect but lacks detail
   0.4 = Somewhat Relevant: Partial related information
   0.5 = Moderately Relevant: Limited relevance
   0.6 = Fairly Relevant: Relevant but lacking depth
   0.7 = Relevant: Clearly relates with substantive information
   0.8 = Very Relevant: Strongly relates with significant information
   0.9 = Highly Relevant: Almost completely answers with detailed information
   1 = Perfectly Relevant: Directly and comprehensively answers the query

3. Additional Guidance:
   - Evaluate based only on content relative to the query
   - Be clear and concise in justifications
   - Do not infer information beyond what's explicitly stated
"""
    
    def rerank_documents(self, query: str, documents: List[Dict], 
                        documents_batch_size: int = 3, llm_weight: float = 0.7) -> List[Dict]:
        """Rerank pages using LLM with relevance score adjustment."""
        if not documents:
            return []
        
        doc_batches = [documents[i:i + documents_batch_size] for i in range(0, len(documents), documents_batch_size)]
        vector_weight = 1 - llm_weight
        
        def process_batch(batch):
            texts = [doc['text'] for doc in batch]
            llm_scores = self._rerank_batch(texts, query)
            
            results = []
            for doc, llm_score in zip(batch, llm_scores):
                doc_with_score = doc.copy()
                doc_with_score['llm_score'] = llm_score
                doc_with_score['relevance_score'] = llm_score
                
                # Convert distance to similarity score (0-1 range)
                # Use exponential decay to handle distances > 1.0
                distance = doc.get('distance', 0.5)
                vector_similarity = max(0.0, min(1.0, 1.0 / (1.0 + distance)))
                combined_score = llm_weight * llm_score + vector_weight * vector_similarity
                doc_with_score['combined_score'] = round(combined_score, 4)
                results.append(doc_with_score)
            
            return results
        
        with ThreadPoolExecutor(max_workers=3) as executor:
            batch_results = list(executor.map(process_batch, doc_batches))
        
        all_results = []
        for batch in batch_results:
            all_results.extend(batch)
        
        all_results.sort(key=lambda x: x['combined_score'], reverse=True)
        return all_results
    
    def _rerank_batch(self, texts: List[str], question: str) -> List[float]:
        """Enhanced unified reranking with detailed reasoning and robust parsing."""
        if not texts:
            return []
        
        blocks_text = ""
        for i, text in enumerate(texts, 1):
            max_length = 1200 if len(texts) == 1 else 800
            truncated_text = text[:max_length] + "..." if len(text) > max_length else text
            blocks_text += f"\nBlock {i}:\n{truncated_text}\n"
        
        user_prompt = f"""
Query: {question}

Text Blocks:
{blocks_text}
"""
        
        schema_str = """
{
  "block_rankings": [
    {
      "reasoning": "string",
      "relevance_score": float
    }
  ]
}
"""
        
        full_prompt = f"{self.system_prompt_multiple}\n\nYour response must be a valid JSON object matching this schema:\n{schema_str}\n\nProvide one ranking object for each of the {len(texts)} blocks in order.\n\n{user_prompt}"
        
        try:
            response = self.llm.invoke(full_prompt)
            if isinstance(response, BaseMessage):
                response_content = str(response.content)
            else:
                response_content = str(response)
            
            rankings = self._parse_rankings_response(response_content, len(texts))
            return [ranking.relevance_score for ranking in rankings.block_rankings]
            
        except Exception as e:
            _log.warning(f"Error in reranking batch: {e}")
            return [0.5] * len(texts)
    
    def _parse_rankings_response(self, response_content: str, expected_count: int) -> RetrievalRankingMultipleBlocks:
        """Parse ranking response with comprehensive fallback strategies."""
        try:
            parsed = json.loads(response_content)
            return RetrievalRankingMultipleBlocks(**parsed)
        except (json.JSONDecodeError, ValueError):
            json_patterns = [
                r'```json\s*(\{.*?\})\s*```',
                r'```\s*(\{.*?\})\s*```',
                r'\{.*?\}'
            ]
            
            for pattern in json_patterns:
                json_match = re.search(pattern, response_content, re.DOTALL)
                if json_match:
                    try:
                        extracted = json_match.group(1) if len(json_match.groups()) > 0 else json_match.group(0)
                        parsed = json.loads(extracted)
                        return RetrievalRankingMultipleBlocks(**parsed)
                    except:
                        continue
            
            score_matches = re.findall(r'(?:score|relevance)[:\s]*([0-9.]+)', response_content, re.IGNORECASE)
            scores = [max(0.0, min(1.0, float(match))) for match in score_matches[:expected_count]]
            
            if len(scores) < expected_count:
                number_matches = re.findall(r'[0-9.]+', response_content)
                potential_scores = []
                for match in number_matches:
                    try:
                        score = float(match)
                        if 0.0 <= score <= 1.0:
                            potential_scores.append(score)
                    except ValueError:
                        continue
                
                while len(scores) < expected_count and potential_scores:
                    scores.append(potential_scores.pop(0))
            
            while len(scores) < expected_count:
                scores.append(0.5)
            
            rankings = []
            for i, score in enumerate(scores):
                rankings.append(RetrievalRankingSingleBlock(
                    reasoning=f"Fallback parsing for block {i+1} - extracted score from unstructured response",
                    relevance_score=score
                ))
            
            return RetrievalRankingMultipleBlocks(block_rankings=rankings)

class VectorRetriever:
    """Vector-based document retrieval using embedding model and vector database."""
    
    def __init__(self, vectorstore):
        self.vectorstore = vectorstore
    
    def retrieve(self, query: str, top_k: int = 30) -> List[Dict]:
        """Retrieve chunks using vector similarity search."""
        if not self.vectorstore:
            return []
        
        docs_with_scores = self.vectorstore.similarity_search_with_score(query, k=top_k)
        
        results = []
        for doc, score in docs_with_scores:
            result = {
                'text': doc.page_content,
                'page': doc.metadata.get('page', 0),
                'chunk': doc.metadata.get('chunk', 1),
                'distance': float(score),
                'source_file': doc.metadata.get('source_file', ''),
                'document_type': doc.metadata.get('document_type', 'unknown'),
                'metadata': doc.metadata
            }
            results.append(result)
        
        return results

class HybridRetriever:
    """Complete retrieval system following the five-stage pipeline."""
    
    def __init__(self, vectorstore, parsed_reports: List[Dict]):
        self.vector_retriever = VectorRetriever(vectorstore)
        self.parent_aggregator = ParentPageAggregator(parsed_reports)
        self.reranker = LLMReranker()
        
    def retrieve(
        self, 
        query: str, 
        llm_reranking_sample_size: int = 30,
        documents_batch_size: int = 2,
        top_n: int = 10,
        llm_weight: float = 0.7
    ) -> List[Dict]:
        """Complete retrieval pipeline with vector search, parent aggregation and LLM reranking."""
        chunk_results = self.vector_retriever.retrieve(
            query=query,
            top_k=llm_reranking_sample_size
        )
        
        parent_results = self.parent_aggregator.aggregate_to_parent_pages(chunk_results)
        
        reranked_results = self.reranker.rerank_documents(
            query=query,
            documents=parent_results,
            documents_batch_size=documents_batch_size,
            llm_weight=llm_weight
        )
        
        return reranked_results[:top_n]

def retrieval_node(state: GraphState) -> GraphState:
    """Execute complete retrieval pipeline."""
    print(f"Starting retrieval for question: {state.question[:100]}...")
    
    retrieval_start = time.time()
    
    retriever = HybridRetriever(state.vectorstore, state.parsed_reports)
    
    reranked_results = retriever.retrieve(
        query=state.question,
        llm_reranking_sample_size=30,
        documents_batch_size=2,
        top_n=10,
        llm_weight=0.7
    )
    
    retrieval_time = time.time() - retrieval_start
    print(f"Retrieval completed with {len(reranked_results)} results in {retrieval_time:.3f}s")
    
    final_context = _assemble_context(reranked_results)
    
    retrieved_docs = []
    for result in reranked_results:
        doc = Document(
            page_content=result['text'],
            metadata=result['metadata']
        )
        retrieved_docs.append(doc)
    
    print(f"Retrieved {len(reranked_results)} final documents")
    
    return GraphState(
        docs=state.docs,
        vectorstore=state.vectorstore,
        question=state.question,
        retrieved_docs=retrieved_docs,
        parsed_reports=state.parsed_reports,
        vector_results=[],
        reranked_results=reranked_results,
        final_context=final_context,
        start_time=state.start_time,
        end_time=state.end_time,
        retrieval_time=retrieval_time,
        generation_time=state.generation_time,
        total_time=state.total_time,
        input_tokens=state.input_tokens,
        output_tokens=state.output_tokens,
        throughput_tokens_per_second=state.throughput_tokens_per_second
    )

def _assemble_context(results: List[Dict]) -> str:
    """Format top pages into final context string with source identification."""
    context_parts = []
    
    for i, result in enumerate(results, 1):
        page_number = result.get('page', '?')
        text = result['text']
        source_file = result.get('source_file', 'unknown')
        document_type = result.get('document_type', 'unknown')
        
        source_id = f"{source_file.replace('.', '_')}_page_{page_number}"
        
        header = f"[Source: {source_id}] - {document_type.upper()} file '{source_file}', Page {page_number}"
        context_parts.append(f"{header}\n\"\"\"\n{text}\n\"\"\"")
    
    return "\n\n---\n\n".join(context_parts)

# ============================================================================
# ANSWERING SYSTEM
# ============================================================================

def rag_node(state: GraphState) -> GraphState:
    """Generate structured answers using enhanced RAG system."""
    print("Generating structured answer...")
    
    generation_start = time.time()
    
    llm = ChatOpenAI(model="gpt-4o-mini", temperature=0.3)
    
    enhanced_prompt = RAGAnswerPrompt()
    context = state.final_context
    
    user_message = enhanced_prompt.user_prompt.format(
        context=context,
        question=state.question
    )
    
    full_prompt = f"{enhanced_prompt.system_prompt_with_schema}\n\n{user_message}"
    
    input_tokens = count_tokens(full_prompt)
    
    try:
        response = llm.invoke(full_prompt)
        response_content = response.content if isinstance(response, BaseMessage) else str(response)
        response_str = str(response_content) if not isinstance(response_content, str) else response_content
        
        structured_answer = _parse_json_response(response_str, state.question)
        
        final_answer = str(structured_answer.get('final_answer', response_str))
        
        print(f"Answer confidence: {structured_answer.get('confidence_level', 'unknown')}")
        print(f"Sources used: {len(structured_answer.get('relevant_sources', []))}")
        
    except Exception as e:
        print(f"Warning: Error in structured answer generation: {e}")
        simple_prompt = f"Answer this question based on the context:\n\nContext: {context}\n\nQuestion: {state.question}"
        response = llm.invoke(simple_prompt)
        fallback_response = response.content if isinstance(response, BaseMessage) else str(response)
        final_answer = str(fallback_response)
        structured_answer = {
            "step_by_step_analysis": "Fallback answer due to parsing error",
            "reasoning_summary": "Error in structured processing",
            "relevant_sources": [],
            "confidence_level": "low",
            "final_answer": final_answer
        }
    
    generation_time = time.time() - generation_start
    output_tokens = count_tokens(final_answer)
    total_tokens = input_tokens + output_tokens
    throughput = calculate_throughput(total_tokens, state.retrieval_time + generation_time)
    
    print(f"Generation completed in {generation_time:.3f}s")
    print(f"Tokens: {input_tokens} input + {output_tokens} output = {total_tokens} total")
    print(f"Throughput: {throughput} tokens/second")
    
    return GraphState(
        docs=state.docs,
        vectorstore=state.vectorstore,
        question=state.question,
        retrieved_docs=state.retrieved_docs,
        answer=final_answer,
        parsed_reports=state.parsed_reports,
        vector_results=state.vector_results,
        reranked_results=state.reranked_results,
        final_context=state.final_context,
        structured_answer=structured_answer,
        start_time=state.start_time,
        end_time=state.end_time,
        retrieval_time=state.retrieval_time,
        generation_time=generation_time,
        total_time=state.total_time,  
        input_tokens=input_tokens,
        output_tokens=output_tokens,
        throughput_tokens_per_second=throughput
    )

def _parse_json_response(response_text: str, question: str) -> Dict:
    """Parse JSON response with multiple fallback strategies."""
    try:
        parsed = json.loads(response_text)
        
        required_fields = {
            "step_by_step_analysis": "Analysis not available",
            "reasoning_summary": "Summary not available", 
            "relevant_sources": [],
            "confidence_level": "medium",
            "final_answer": "Answer not available"
        }
        
        for field, default in required_fields.items():
            if field not in parsed:
                parsed[field] = default
        
        return parsed
        
    except (json.JSONDecodeError, ValueError):
        json_patterns = [
            r'```json\s*(\{.*?\})\s*```',
            r'```\s*(\{.*?\})\s*```',
            r'\{.*?\}'
        ]
        
        for pattern in json_patterns:
            json_match = re.search(pattern, response_text, re.DOTALL)
            if json_match:
                try:
                    extracted = json_match.group(1) if len(json_match.groups()) > 0 else json_match.group(0)
                    parsed = json.loads(extracted)
                    
                    required_fields = {
                        "step_by_step_analysis": "Analysis extracted from response",
                        "reasoning_summary": "Extracted response", 
                        "relevant_sources": [],
                        "confidence_level": "low",
                        "final_answer": parsed.get("final_answer", response_text[:500])
                    }
                    
                    for field, default in required_fields.items():
                        if field not in parsed:
                            parsed[field] = default
                    
                    return parsed
                except:
                    continue
        
        return {
            "step_by_step_analysis": f"Unable to parse structured analysis. Raw response: {response_text[:500]}...",
            "reasoning_summary": "Fallback parsing used due to malformed JSON",
            "relevant_sources": [],
            "confidence_level": "low",
            "final_answer": response_text[:1000]
        }

def log_node(state: GraphState) -> GraphState:
    """Log detailed metrics and results including performance data."""
    print("Logging enhanced metrics...")
    
    end_time = time.time()
    total_time = end_time - state.start_time if state.start_time > 0 else 0.0
    
    total_tokens = state.input_tokens + state.output_tokens
    precise_throughput = calculate_throughput(total_tokens, total_time)
    
    metrics = {}
    if state.reranked_results:
        scores = [r.get('combined_score', 0) for r in state.reranked_results]
        if scores:
            metrics = {
                "avg_score": round(sum(scores) / len(scores), 4),
                "max_score": round(max(scores), 4),
                "score_range": round(max(scores) - min(scores), 4),
                "final_results": len(state.reranked_results)
            }
    
    performance_metrics = {
        "retrieval_time_s": round(state.retrieval_time, 3),
        "generation_time_s": round(state.generation_time, 3),
        "total_time_s": round(total_time, 3),  
        "input_tokens": state.input_tokens,
        "output_tokens": state.output_tokens,
        "total_tokens": total_tokens,
        "throughput_tokens_per_second": precise_throughput,  
        "start_time": state.start_time,
        "end_time": end_time
    }
    
    log_entry = {
        "question": state.question,
        "answer": state.answer,
        "confidence_level": state.structured_answer.get('confidence_level', 'unknown'),
        "relevant_sources": state.structured_answer.get('relevant_sources', []),
        "reasoning_summary": state.structured_answer.get('reasoning_summary', ''),
        "retrieval_metrics": metrics,
        "performance_metrics": performance_metrics,
        "used_existing_vectordb": False
    }
    
    df = pd.DataFrame([log_entry])
    log_file = "enhanced_rag_qa_log.csv"
    
    if not os.path.exists(log_file):
        df.to_csv(log_file, index=False)
        print(f"Created new log file: {log_file}")
    else:
        df.to_csv(log_file, mode="a", header=False, index=False)
        print(f"Appended to log file: {log_file}")
    
    print("\n" + "=" * 50)
    print("PERFORMANCE SUMMARY")
    print("=" * 50)
    print(f"Retrieval Time: {performance_metrics['retrieval_time_s']}s")
    print(f"Generation Time: {performance_metrics['generation_time_s']}s")
    print(f"Total Processing Time: {performance_metrics['total_time_s']}s")
    print(f"Input Tokens: {performance_metrics['input_tokens']:,}")
    print(f"Output Tokens: {performance_metrics['output_tokens']:,}")
    print(f"Total Tokens: {performance_metrics['total_tokens']:,}")
    print(f"Throughput: {performance_metrics['throughput_tokens_per_second']} tokens/second")
    print("=" * 50)
    
    updated_state = GraphState(
        docs=state.docs,
        vectorstore=state.vectorstore,
        question=state.question,
        retrieved_docs=state.retrieved_docs,
        answer=state.answer,
        parsed_reports=state.parsed_reports,
        vector_results=state.vector_results,
        reranked_results=state.reranked_results,
        final_context=state.final_context,
        structured_answer=state.structured_answer,
        start_time=state.start_time,
        end_time=end_time,
        retrieval_time=state.retrieval_time,
        generation_time=state.generation_time,
        total_time=total_time,
        input_tokens=state.input_tokens,
        output_tokens=state.output_tokens,
        throughput_tokens_per_second=precise_throughput
    )
    
    return updated_state

# ============================================================================
# WORKFLOW ORCHESTRATION
# ============================================================================

init_workflow = StateGraph(GraphState)
init_workflow.add_node("ingest", ingest_node)
init_workflow.add_node("embed", embed_node)
init_workflow.set_entry_point("ingest")
init_workflow.add_edge("ingest", "embed")
init_workflow.add_edge("embed", END)
init_graph = init_workflow.compile()

query_workflow = StateGraph(GraphState)
query_workflow.add_node("retrieval", retrieval_node)
query_workflow.add_node("rag", rag_node)
query_workflow.add_node("log", log_node)
query_workflow.set_entry_point("retrieval")
query_workflow.add_edge("retrieval", "rag")
query_workflow.add_edge("rag", "log")
query_workflow.add_edge("log", END)
query_graph = query_workflow.compile()

# ============================================================================
# MAIN EXECUTION
# ============================================================================

if __name__ == "__main__":
    print("RAG system initialization...")
    
    vs_manager = VectorStoreManager()
    
    if vs_manager.vectorstore_exists():
        try:
            vectorstore = vs_manager.load_existing_vectorstore()
            if vectorstore is not None:
                stats = vs_manager.get_vectorstore_stats(vectorstore)
                print(f"Found existing vector database: {stats}")
                docs = []
                parsed_reports = []
            else:
                raise Exception("Failed to load existing database")
        except Exception as e:
            print(f"Error loading existing database: {e}")
            print("Creating new database from documents...")
            vectorstore = None
    else:
        print("No existing vector database found.")
        vectorstore = None
    
    if vectorstore is None:
        print("Please provide documents for processing:")
        user_files = get_user_files()
        
        if not user_files:
            print("No files provided. Exiting.")
            exit(1)
        
        print(f"Processing {len(user_files)} documents...")
        
        initial_state = GraphState(
            docs=user_files,
            question="",
            skip_parsing=False
        )
        
        try:
            state_after_embed = init_graph.invoke(initial_state)
            print("Documents processed successfully!")
            
            docs = state_after_embed["docs"]
            vectorstore = state_after_embed["vectorstore"]
            parsed_reports = state_after_embed["parsed_reports"]
            
            if parsed_reports:
                total_pages = sum(len(report['report']['content']['pages']) for report in parsed_reports)
                total_chunks = len(docs)
                print(f"Documents: {len(parsed_reports)}, Pages: {total_pages}, Chunks: {total_chunks}")
            
        except Exception as e:
            print(f"Error processing documents: {e}")
            print("Check: dependencies, file formats, file paths")
            exit(1)
    
    print("==== Ready for questions. Type 'quit' to exit. ====")
    
    while True:
        question = input("\nQuestion: ")
        
        if question.strip().lower() == "quit":
            print("Exiting...")
            break
        
        start_time = time.time()
        
        state = GraphState(
            docs=docs,
            vectorstore=vectorstore,
            question=question,
            parsed_reports=parsed_reports,
            start_time=start_time
        )
        
        result = query_graph.invoke(state)
        
        print("\n" + "=" * 60)
        print("RESULTS")
        print("=" * 60)
        print(f"Confidence: {result['structured_answer'].get('confidence_level', 'unknown')}")
        print(f"Documents: {len(result.get('reranked_results', []))}")
        
        if result.get('reranked_results'):
            scores = [r.get('combined_score', 0) for r in result['reranked_results']]
            if scores:
                print(f"Score Range: {min(scores):.3f} - {max(scores):.3f}")
        
        print("\n--- Reasoning ---")
        print(result['structured_answer'].get('reasoning_summary', 'No summary'))
        
        print("\n--- Analysis ---")
        analysis = result['structured_answer'].get('step_by_step_analysis', 'No analysis')
        print(analysis[:400] + "..." if len(analysis) > 400 else analysis)
        
        print("\n--- Sources ---")
        sources = result['structured_answer'].get('relevant_sources', [])
        if sources:
            for i, source in enumerate(sources, 1):
                print(f"  {i}. {source}")
        else:
            print("  No sources identified")
        
        print("\n" + "=" * 60)
        print("ANSWER")
        print("=" * 60)
        print(result["answer"])
        print("=" * 60)
